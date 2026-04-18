# Unified document extractor: routes each file type to the correct extraction pipeline.
# Returns list[PageResult] compatible with existing embedder.embed_pages().

from __future__ import annotations

import asyncio
from pathlib import Path

from agents.audit.processors.pdf_extractor import PageResult, extract_pages
from core.logger import get_logger

logger = get_logger("audit.document_extractor")

SUPPORTED_EXTENSIONS = frozenset({
    ".pdf", ".pptx", ".ppt", ".docx", ".doc", ".xlsx", ".csv"
})

_SUPPORTED_LABEL = "PDF, PPTX, PPT, DOCX, DOC, XLSX, CSV"


# ---------------------------------------------------------------------------
# Per-format extractors
# ---------------------------------------------------------------------------

def _extract_pdf(file_path: str) -> list[PageResult]:
    with open(file_path, "rb") as fh:
        pdf_bytes = fh.read()
    return extract_pages(pdf_bytes)


def _extract_pptx(file_path: str, filename: str) -> list[PageResult]:
    try:
        from pptx import Presentation  # type: ignore
    except ImportError as exc:
        raise RuntimeError("python-pptx is required. Run: pip install python-pptx") from exc

    try:
        prs = Presentation(file_path)
    except Exception as exc:
        raise ValueError(
            f"Could not open '{filename}' as a PowerPoint file. "
            "Note: legacy .ppt binary format is not supported — please convert to .pptx first."
        ) from exc

    results: list[PageResult] = []

    for slide_num, slide in enumerate(prs.slides, start=1):
        parts: list[str] = []

        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = " ".join(run.text for run in para.runs if run.text.strip())
                    if line.strip():
                        parts.append(line.strip())

            if shape.has_table:
                rows: list[str] = []
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    rows.append(" | ".join(cells))
                if rows:
                    parts.append("\n".join(rows))

        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                parts.append(f"[Speaker Notes] {notes}")

        text = "\n".join(parts).strip()
        results.append(
            PageResult(
                page_num=slide_num,
                text=text,
                tables=[],
                image_bytes=b"",
                needs_vision=False,
            )
        )
        logger.info(f"slide_extracted filename={filename} slide={slide_num} chars={len(text)}")

    return results


def _extract_docx(file_path: str, filename: str) -> list[PageResult]:
    try:
        from docx import Document  # type: ignore
    except ImportError as exc:
        raise RuntimeError("python-docx is required. Run: pip install python-docx") from exc

    try:
        doc = Document(file_path)
    except Exception as exc:
        raise ValueError(
            f"Could not open '{filename}' as a Word document. "
            "Note: legacy .doc binary format is not supported — please convert to .docx first."
        ) from exc

    PARAS_PER_PAGE = 30
    all_parts: list[str] = []

    for para in doc.paragraphs:
        text = para.text.strip()
        if text:
            all_parts.append(text)

    for table in doc.tables:
        rows: list[str] = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            rows.append(" | ".join(cells))
        if rows:
            all_parts.append("\n".join(rows))

    if not all_parts:
        return [PageResult(page_num=1, text="", tables=[], image_bytes=b"", needs_vision=False)]

    results: list[PageResult] = []
    for i in range(0, len(all_parts), PARAS_PER_PAGE):
        page_num = i // PARAS_PER_PAGE + 1
        text = "\n".join(all_parts[i : i + PARAS_PER_PAGE]).strip()
        results.append(
            PageResult(
                page_num=page_num,
                text=text,
                tables=[],
                image_bytes=b"",
                needs_vision=False,
            )
        )
        logger.info(f"docx_page_extracted filename={filename} page={page_num} chars={len(text)}")

    return results


def _extract_xlsx(file_path: str, filename: str) -> list[PageResult]:
    try:
        import openpyxl  # already in requirements
    except ImportError as exc:
        raise RuntimeError("openpyxl is required. Run: pip install openpyxl") from exc

    wb = openpyxl.load_workbook(file_path, data_only=True)
    results: list[PageResult] = []

    for sheet_num, sheet_name in enumerate(wb.sheetnames, start=1):
        ws = wb[sheet_name]
        lines: list[str] = [f"Sheet: {sheet_name}"]

        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            if any(c.strip() for c in cells):
                lines.append(" | ".join(cells).strip(" |"))

        text = "\n".join(lines).strip()
        results.append(
            PageResult(
                page_num=sheet_num,
                text=text,
                tables=[],
                image_bytes=b"",
                needs_vision=False,
            )
        )
        logger.info(f"xlsx_sheet_extracted filename={filename} sheet={sheet_name} chars={len(text)}")

    return results or [PageResult(page_num=1, text="", tables=[], image_bytes=b"", needs_vision=False)]


def _extract_csv(file_path: str, filename: str) -> list[PageResult]:
    try:
        import pandas as pd  # already in requirements
    except ImportError as exc:
        raise RuntimeError("pandas is required. Run: pip install pandas") from exc

    df = pd.read_csv(file_path)
    ROWS_PER_PAGE = 100
    results: list[PageResult] = []
    header = " | ".join(str(c) for c in df.columns)

    for i in range(0, max(1, len(df)), ROWS_PER_PAGE):
        page_num = i // ROWS_PER_PAGE + 1
        chunk = df.iloc[i : i + ROWS_PER_PAGE]
        lines = [header]
        for _, row in chunk.iterrows():
            lines.append(" | ".join(str(v) for v in row.values))
        text = "\n".join(lines).strip()
        results.append(
            PageResult(
                page_num=page_num,
                text=text,
                tables=[],
                image_bytes=b"",
                needs_vision=False,
            )
        )
        logger.info(f"csv_page_extracted filename={filename} page={page_num} rows={len(chunk)}")

    return results or [PageResult(page_num=1, text="", tables=[], image_bytes=b"", needs_vision=False)]


# ---------------------------------------------------------------------------
# Public router function
# ---------------------------------------------------------------------------

def extract_document(file_path: str, filename: str) -> list[PageResult]:
    """
    Detect file type from extension and extract content as a list of PageResult.

    Supported extensions: PDF, PPTX, PPT, DOCX, DOC, XLSX, CSV.
    For non-PDF files, image_bytes=b"" and needs_vision=False on every page.
    Each slide / sheet / paragraph-group is treated as a single "page".

    Raises ValueError for unsupported extensions.
    """
    ext = Path(filename).suffix.lower()

    if ext not in SUPPORTED_EXTENSIONS:
        label = ext.lstrip(".").upper() if ext else "unknown"
        raise ValueError(
            f"File type '{label}' is not supported. Supported types: {_SUPPORTED_LABEL}"
        )

    logger.info(f"extract_document_start filename={filename} ext={ext}")

    if ext == ".pdf":
        return _extract_pdf(file_path)
    if ext in (".pptx", ".ppt"):
        return _extract_pptx(file_path, filename)
    if ext in (".docx", ".doc"):
        return _extract_docx(file_path, filename)
    if ext == ".xlsx":
        return _extract_xlsx(file_path, filename)
    # .csv
    return _extract_csv(file_path, filename)
